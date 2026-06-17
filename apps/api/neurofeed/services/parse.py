"""Document parsing. Dispatch by source type → ordered text blocks with page/slide refs.

Each block is a dict: {"text": str, "page_ref": {"page": int} | {"slide": int} | {}, "heading": str | None}
The chunker downstream stitches blocks within a page/slide and splits across the token budget.
"""
from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Literal


SourceType = Literal["pdf", "docx", "pptx", "audio", "text"]


@dataclass
class Block:
    text: str
    page_ref: dict
    heading: str | None = None


def parse(source_type: SourceType, data: bytes, *, filename: str | None = None) -> list[Block]:
    if source_type == "pdf":
        return _parse_pdf(data)
    if source_type == "docx":
        return _parse_docx(data)
    if source_type == "pptx":
        return _parse_pptx(data)
    if source_type == "audio":
        return _parse_audio(data, filename=filename)
    if source_type == "text":
        return [Block(text=data.decode("utf-8", errors="replace"), page_ref={})]
    raise ValueError(f"unsupported source_type: {source_type}")


# ---------- PDF ----------
def _parse_pdf(data: bytes) -> list[Block]:
    import fitz  # PyMuPDF

    blocks: list[Block] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page_index in range(len(doc)):
            page = doc[page_index]
            # get_text("blocks") returns (x0, y0, x1, y1, text, block_no, block_type)
            raw_blocks = page.get_text("blocks") or []
            raw_blocks.sort(key=lambda b: (round(b[1], 1), round(b[0], 1)))
            for b in raw_blocks:
                text = (b[4] or "").strip()
                if not text:
                    continue
                heading = _detect_heading(text)
                blocks.append(Block(text=text, page_ref={"page": page_index + 1}, heading=heading))
    return blocks


# ---------- DOCX ----------
def _parse_docx(data: bytes) -> list[Block]:
    from docx import Document  # python-docx

    blocks: list[Block] = []
    doc = Document(BytesIO(data))
    current_heading: str | None = None
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        is_heading = style.lower().startswith("heading") or _detect_heading(text) is not None
        if is_heading:
            current_heading = text
        blocks.append(Block(text=text, page_ref={}, heading=current_heading if not is_heading else text))
    return blocks


# ---------- PPTX ----------
def _parse_pptx(data: bytes) -> list[Block]:
    from pptx import Presentation

    blocks: list[Block] = []
    prs = Presentation(BytesIO(data))
    for slide_index, slide in enumerate(prs.slides):
        slide_title: str | None = None
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip() or None
        except Exception:
            pass
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                blocks.append(
                    Block(
                        text=text,
                        page_ref={"slide": slide_index + 1},
                        heading=slide_title,
                    )
                )
    return blocks


# ---------- Audio (Groq Whisper) ----------
def _parse_audio(data: bytes, *, filename: str | None) -> list[Block]:
    """Synchronous wrapper around Groq Whisper. Returns a single block per segment."""
    from ..config import get_settings
    import httpx

    s = get_settings()
    if not s.groq_api_key:
        raise RuntimeError("GROQ_API_KEY required for audio transcription")

    files = {
        "file": (filename or "audio.bin", data, "application/octet-stream"),
        "model": (None, s.groq_stt_model),
        "response_format": (None, "verbose_json"),
    }
    headers = {"Authorization": f"Bearer {s.groq_api_key}"}
    url = s.groq_base_url.rstrip("/") + "/audio/transcriptions"
    with httpx.Client(timeout=180.0) as client:
        r = client.post(url, files=files, headers=headers)
        r.raise_for_status()
        payload = r.json()

    segments = payload.get("segments") or []
    if not segments:
        text = (payload.get("text") or "").strip()
        return [Block(text=text, page_ref={})] if text else []

    blocks: list[Block] = []
    for i, seg in enumerate(segments):
        text = (seg.get("text") or "").strip()
        if not text:
            continue
        blocks.append(
            Block(
                text=text,
                page_ref={"segment": i + 1, "t_start": seg.get("start"), "t_end": seg.get("end")},
                heading=None,
            )
        )
    return blocks


# ---------- Helpers ----------
def _detect_heading(text: str) -> str | None:
    """Very cheap heading detector: short, title-cased, no terminal punctuation."""
    if len(text) > 100 or "\n" in text:
        return None
    if text.endswith((".", "?", "!", ":")):
        return None
    words = text.split()
    if not words or len(words) > 12:
        return None
    title_like = sum(1 for w in words if w[:1].isupper())
    if title_like / len(words) >= 0.6:
        return text
    return None


def detect_source_type(filename: str, content_type: str | None) -> SourceType:
    name = filename.lower()
    if name.endswith(".pdf"):
        return "pdf"
    if name.endswith(".docx"):
        return "docx"
    if name.endswith(".pptx"):
        return "pptx"
    if name.endswith((".mp3", ".m4a", ".wav", ".ogg", ".webm", ".flac")):
        return "audio"
    if name.endswith((".txt", ".md")):
        return "text"
    ct = (content_type or "").lower()
    if "pdf" in ct:
        return "pdf"
    if "word" in ct or "officedocument.wordprocessingml" in ct:
        return "docx"
    if "presentation" in ct or "officedocument.presentationml" in ct:
        return "pptx"
    if ct.startswith("audio/"):
        return "audio"
    if ct.startswith("text/"):
        return "text"
    raise ValueError(f"cannot infer source type from filename={filename} content_type={content_type}")
